"""
PowerPoint (PPTX) text extractor using python-pptx.
Extracts title, text frames, tables, and speaker notes from each slide.
"""
import logging
from typing import List, Dict

from pptx import Presentation

logger = logging.getLogger(__name__)


def extract_pptx_text(file_path: str) -> List[Dict]:
    """
    Extract text from each slide of a PPTX file.
    Returns: [{page_number (slide number), text}]
    """
    prs = Presentation(file_path)
    pages: List[Dict] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []

        # Title
        if slide.shapes.title and slide.shapes.title.text.strip():
            parts.append(slide.shapes.title.text.strip())

        # All text frames
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(
                        run.text for run in para.runs if run.text.strip()
                    ).strip()
                    if line:
                        parts.append(line)

            # Tables
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        parts.append(row_text)

        # Speaker notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame and notes_frame.text.strip():
                parts.append(notes_frame.text.strip())

        slide_text = "\n".join(parts).strip()
        if slide_text:
            pages.append({"page_number": slide_num, "text": slide_text})
            logger.debug(f"  Slide {slide_num}: {len(slide_text)} chars")

    logger.info(f"PPTX extracted: {len(pages)} slides with content")
    return pages
